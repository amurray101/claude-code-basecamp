#!/usr/bin/env python3
"""
Build the Claude Code Basecamp facilitator slide deck.
Run: python3 build_slides.py
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ──────────────────────────────────────────────────────────────
BG       = RGBColor(0xFA, 0xF9, 0xF5)  # warm cream
TEXT     = RGBColor(0x14, 0x14, 0x13)  # near-black
MUTED    = RGBColor(0x6A, 0x68, 0x5E)
FAINT    = RGBColor(0xB0, 0xAE, 0xA5)
ORANGE   = RGBColor(0xD9, 0x77, 0x57)
BLUE     = RGBColor(0x6A, 0x9B, 0xCC)
GREEN    = RGBColor(0x78, 0x8C, 0x5D)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FOOTER_TEXT = "CLAUDE CODE BASECAMP  |  Becoming an Anthropic Pro"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use a blank layout
blank_layout = prs.slide_layouts[6]  # blank

slide_number = [0]  # mutable counter


# ── Helper functions ─────────────────────────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=TEXT, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf


def add_bullet_slide_tf(slide, left, top, width, height, items,
                         font_size=18, color=TEXT, font_name="Calibri",
                         bold_first=False, line_spacing=1.35):
    """Add a text frame with bullet points. items is a list of strings."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        p.line_spacing = Pt(font_size * line_spacing)
        if bold_first:
            # bold everything before the first em-dash or colon
            pass  # keep simple
    return tf


def add_footer(slide):
    slide_number[0] += 1
    # Footer bar
    add_textbox(slide, Inches(0.75), Inches(6.9), Inches(10), Inches(0.4),
                FOOTER_TEXT, font_size=9, color=FAINT,
                alignment=PP_ALIGN.LEFT, line_spacing=1.0)
    # Slide number
    add_textbox(slide, Inches(11.8), Inches(6.9), Inches(1), Inches(0.4),
                str(slide_number[0]), font_size=9, color=FAINT,
                alignment=PP_ALIGN.RIGHT, line_spacing=1.0)


def add_accent_bar(slide, color, left=Inches(0.75), top=Inches(1.6),
                   width=Inches(1.2), height=Inches(0.06)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def make_content_slide(title, bullets, accent_color=MUTED, subtitle=None):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, BG)

    # Title
    add_textbox(slide, Inches(0.75), Inches(0.5), Inches(11), Inches(0.9),
                title, font_size=32, bold=True, color=TEXT)

    # Accent bar under title
    add_accent_bar(slide, accent_color, top=Inches(1.35))

    # Optional subtitle
    y_start = Inches(1.7)
    if subtitle:
        add_textbox(slide, Inches(0.75), y_start, Inches(11), Inches(0.5),
                    subtitle, font_size=16, color=MUTED)
        y_start = Inches(2.2)

    # Bullets
    if bullets:
        add_bullet_slide_tf(slide, Inches(0.75), y_start, Inches(11), Inches(4.5),
                            bullets, font_size=19, color=TEXT)

    add_footer(slide)
    return slide


def make_two_column_slide(title, left_title, left_items, right_title, right_items,
                          accent_color=MUTED):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, BG)

    add_textbox(slide, Inches(0.75), Inches(0.5), Inches(11), Inches(0.9),
                title, font_size=32, bold=True, color=TEXT)
    add_accent_bar(slide, accent_color, top=Inches(1.35))

    # Left column header
    add_textbox(slide, Inches(0.75), Inches(1.8), Inches(5.5), Inches(0.5),
                left_title, font_size=20, bold=True, color=accent_color)
    add_bullet_slide_tf(slide, Inches(0.75), Inches(2.35), Inches(5.5), Inches(4),
                        left_items, font_size=17, color=TEXT)

    # Divider line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(6.55), Inches(1.8), Inches(0.02), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = FAINT
    shape.line.fill.background()

    # Right column header
    add_textbox(slide, Inches(6.9), Inches(1.8), Inches(5.5), Inches(0.5),
                right_title, font_size=20, bold=True, color=accent_color)
    add_bullet_slide_tf(slide, Inches(6.9), Inches(2.35), Inches(5.5), Inches(4),
                        right_items, font_size=17, color=TEXT)

    add_footer(slide)
    return slide


def make_section_divider(day_label, day_title, day_subtitle, accent_color):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, BG)

    # Large accent bar at top
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(0), SLIDE_W, Inches(0.18))
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()

    # Day label
    add_textbox(slide, Inches(0.75), Inches(2.0), Inches(11), Inches(0.7),
                day_label, font_size=18, bold=True, color=accent_color,
                font_name="Calibri")

    # Day title
    add_textbox(slide, Inches(0.75), Inches(2.7), Inches(11), Inches(1.0),
                day_title, font_size=42, bold=True, color=TEXT)

    # Subtitle
    add_textbox(slide, Inches(0.75), Inches(3.8), Inches(11), Inches(0.6),
                day_subtitle, font_size=20, color=MUTED)

    # Accent bar
    add_accent_bar(slide, accent_color, top=Inches(4.6), width=Inches(2.5),
                   height=Inches(0.05))

    add_footer(slide)
    return slide


# ═════════════════════════════════════════════════════════════════════════
# TITLE SLIDE
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)

# Accent stripe at very top
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0), Inches(0), SLIDE_W, Inches(0.12))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

add_textbox(slide, Inches(0.75), Inches(1.8), Inches(11), Inches(1.2),
            "Claude Code Basecamp", font_size=48, bold=True, color=TEXT)
add_textbox(slide, Inches(0.75), Inches(3.1), Inches(11), Inches(0.6),
            "Five days to technical fluency", font_size=24, color=MUTED)
add_textbox(slide, Inches(0.75), Inches(3.8), Inches(11), Inches(0.5),
            "Facilitator Deck", font_size=18, color=FAINT)

add_accent_bar(slide, ORANGE, top=Inches(4.6), width=Inches(2.0), height=Inches(0.05))

add_textbox(slide, Inches(0.75), Inches(5.2), Inches(11), Inches(0.5),
            "Taught by [Facilitator Name]", font_size=16, color=MUTED)

add_footer(slide)


# ═════════════════════════════════════════════════════════════════════════
# DAY 1: FIRST CONTACT (orange)
# ═════════════════════════════════════════════════════════════════════════
make_section_divider("DAY 1", "First Contact",
                     "Getting started with Claude Code", ORANGE)

# 1.1 The Agentic Difference
make_two_column_slide(
    "The Agentic Difference",
    "Autocomplete Tools",
    [
        "Suggest code line-by-line",
        "You accept or reject each suggestion",
        "Reactive — waits for your cursor",
        "Works within a single file",
    ],
    "Agentic Coding",
    [
        "Describe the outcome you want",
        "Claude plans, then executes across files",
        "Proactive — reads your whole codebase",
        "Runs tests, iterates, self-corrects",
    ],
    accent_color=ORANGE,
)
# Key message sub-slide
s = prs.slides[-1]
add_textbox(s, Inches(0.75), Inches(5.8), Inches(11.5), Inches(0.8),
            '"Claude Code reads your codebase, plans changes, edits multiple files, '
            'runs tests — it\'s a colleague, not a clipboard."',
            font_size=16, color=ORANGE, alignment=PP_ALIGN.LEFT)

# 1.2 The Agentic Loop
make_content_slide(
    "The Agentic Loop",
    [
        "1.  Read context — package.json, directory structure, CLAUDE.md",
        "2.  Plan approach — decide which files to touch, what strategy to use",
        "3.  Execute changes — edit files, create new ones, update imports",
        "4.  Verify results — run tests, lint, type-check",
        "5.  Iterate — fix failures, refine until green",
        "",
        "Claude reads your project BEFORE you type anything.",
        "Every response is grounded in your actual code.",
    ],
    accent_color=ORANGE,
)

# 1.3 Two Surfaces, One Engine
make_two_column_slide(
    "Two Surfaces, One Engine",
    "CLI (Terminal)",
    [
        "Raw speed — launch and go",
        "Scriptable — pipe inputs, chain commands",
        "CI/CD integration — automate workflows",
        "Perfect for power users",
    ],
    "IDE (VS Code / JetBrains)",
    [
        "Inline diffs — see changes in context",
        "Visual file tree — navigate with clicks",
        "Familiar UX — feels like your editor",
        "Perfect for visual thinkers",
    ],
    accent_color=ORANGE,
)
s = prs.slides[-1]
add_textbox(s, Inches(0.75), Inches(5.8), Inches(11.5), Inches(0.6),
            "Same agentic engine underneath. Choose the surface that fits your style.",
            font_size=16, color=ORANGE)

# 1.4 Live Demo: Install + First Task
make_content_slide(
    "Live Demo: Install + First Task",
    [
        "1.  Install the CLI:  npm install -g @anthropic-ai/claude-code",
        "2.  Authenticate:  claude  (follow browser prompt)",
        "3.  Clone sample repo, cd into it",
        "4.  Launch Claude Code:  claude",
        "5.  Give it a real task:",
        '     "Add a /health endpoint that returns { status: ok, uptime: ... }.',
        '      Include unit tests."',
        "6.  Watch the agentic loop in action",
    ],
    accent_color=ORANGE,
    subtitle="Facilitator demo script",
)

# 1.5 Discussion: What Did You Notice?
make_content_slide(
    "Discussion: What Did You Notice?",
    [
        "What files did Claude read first?",
        "",
        "How did it decide what to change?",
        "",
        "How did it verify its own work?",
        "",
        "How does this compare to Copilot or ChatGPT?",
        "",
        "What surprised you?",
    ],
    accent_color=ORANGE,
    subtitle="Group discussion — 10 minutes",
)


# ═════════════════════════════════════════════════════════════════════════
# DAY 2: PROMPT CRAFT (blue)
# ═════════════════════════════════════════════════════════════════════════
make_section_divider("DAY 2", "Prompt Craft",
                     "CLAUDE.md, context management, and the art of steering", BLUE)

# 2.1 Before & After: The CLAUDE.md Effect
make_two_column_slide(
    "Before & After: The CLAUDE.md Effect",
    "Without CLAUDE.md",
    [
        "Claude guesses your conventions",
        "Inconsistent naming, style, patterns",
        "May pick wrong test framework",
        "Ignores team-specific rules",
        "Output feels generic",
    ],
    "With CLAUDE.md",
    [
        "Claude follows your patterns exactly",
        "Consistent naming, style, structure",
        "Uses your test framework + conventions",
        "Respects team-specific rules",
        "Output feels like a team member wrote it",
    ],
    accent_color=BLUE,
)
s = prs.slides[-1]
add_textbox(s, Inches(0.75), Inches(5.8), Inches(11.5), Inches(0.6),
            "CLAUDE.md is the single most important configuration surface.",
            font_size=16, color=BLUE, bold=True)

# 2.2 Anatomy of a CLAUDE.md
make_content_slide(
    "Anatomy of a CLAUDE.md",
    [
        "# Architecture — What does the system look like? Key modules, data flow.",
        "",
        "# Conventions — Naming, file structure, import order, error handling.",
        "",
        "# Testing — Framework, patterns, where tests live, how to run them.",
        "",
        "# Before Committing — Lint, type-check, test commands to run.",
        "",
        "Each section grounds Claude in YOUR team's way of working.",
        "Think of it as onboarding docs — for your AI colleague.",
    ],
    accent_color=BLUE,
)

# 2.3 Live Demo: Writing a CLAUDE.md
make_content_slide(
    "Live Demo: Writing a CLAUDE.md",
    [
        "1.  Open a messy repo (no CLAUDE.md yet)",
        "2.  Ask Claude to make a change — observe inconsistencies",
        "3.  Write a CLAUDE.md together:",
        '     Architecture, Conventions, Testing, Before Committing',
        "4.  Ask Claude the same task again",
        "5.  Compare output quality — before vs. after",
        "",
        "The difference is dramatic. Show it live.",
    ],
    accent_color=BLUE,
    subtitle="Facilitator demo script",
)

# 2.4 Session Management
make_content_slide(
    "Session Management",
    [
        "/compact — Summarize conversation to save context window",
        "  Use when: session is getting long, Claude starts repeating itself",
        "",
        "/clear — Reset the conversation completely",
        "  Use when: switching to a new task, context is polluted",
        "",
        "/cost — Check token usage and spend",
        "  Use when: monitoring costs, before/after comparisons",
        "",
        "Long sessions degrade quality. Compact early, compact often.",
    ],
    accent_color=BLUE,
    subtitle="Essential commands for staying effective",
)

# 2.5 Plan Mode
make_content_slide(
    "Plan Mode",
    [
        'Prefix any prompt with "plan:" to think before acting.',
        "",
        "Claude will outline its approach WITHOUT making changes.",
        "",
        "Perfect for:",
        "  Complex, multi-file refactors",
        "  Customer demos (shows reasoning without risk)",
        "  Building trust with skeptical engineers",
        "  Reviewing strategy before committing to execution",
        "",
        "\"plan: How would you migrate this Express app to Fastify?\"",
    ],
    accent_color=BLUE,
)

# 2.6 Prompt Patterns That Work
make_two_column_slide(
    "Prompt Patterns That Work",
    "Too Vague",
    [
        '"Refactor src/utils"',
        '"Fix the tests"',
        '"Make it faster"',
        '"Clean this up"',
    ],
    "Specific & Actionable",
    [
        '"Refactor src/utils/helpers.js to use async/await,',
        '  add JSDoc, write co-located tests per CLAUDE.md"',
        '"Fix the failing test in auth.test.ts — the mock',
        '  is returning undefined instead of a user object"',
    ],
    accent_color=BLUE,
)
s = prs.slides[-1]
add_textbox(s, Inches(0.75), Inches(5.8), Inches(11.5), Inches(0.6),
            "Specificity is kindness. The more context you give, the better the output.",
            font_size=16, color=BLUE)

# 2.7 Common Anti-Patterns
make_content_slide(
    "Common Anti-Patterns",
    [
        "Over-prompting — Giving Claude 15 constraints at once. It loses focus.",
        "",
        "Under-specifying — \"Make it better\" gives Claude nothing to work with.",
        "",
        "Ignoring CLAUDE.md — Your best leverage, unused.",
        "",
        "Never compacting — Long sessions = degraded quality. Use /compact.",
        "",
        "Copy-pasting code into the prompt — Claude already has your files.",
        "  Just reference them by path.",
    ],
    accent_color=BLUE,
)

# 2.8 Discussion: Your Team's Conventions
make_content_slide(
    "Discussion: Your Team's Conventions",
    [
        "What conventions would you codify in a CLAUDE.md?",
        "",
        "What patterns does your team argue about most?",
        "",
        "How does CLAUDE.md change the customer conversation?",
        "  (Hint: it's not just for developers — it's a governance tool)",
        "",
        "What would you put in your customer's first CLAUDE.md?",
    ],
    accent_color=BLUE,
    subtitle="Group discussion — 10 minutes",
)


# ═════════════════════════════════════════════════════════════════════════
# DAY 3: EXTEND & CUSTOMIZE (green)
# ═════════════════════════════════════════════════════════════════════════
make_section_divider("DAY 3", "Extend & Customize",
                     "Hooks, MCP, and the integration architecture", GREEN)

# 3.1 The Integration Layer
make_content_slide(
    "The Integration Layer",
    [
        "Three extension points that make Claude Code enterprise-ready:",
        "",
        "HOOKS  (Guardrails)",
        "  Control what Claude does — enforce quality gates automatically",
        "",
        "MCP  (Connections)",
        "  Connect Claude to external tools — Jira, Slack, Datadog, APIs",
        "",
        "SLASH COMMANDS  (Workflows)",
        "  Team-shareable procedures — define once, everyone inherits",
        "",
        "Together: guardrails + connections + workflows = enterprise platform.",
    ],
    accent_color=GREEN,
)

# 3.2 Hooks: Automated Guardrails
make_content_slide(
    "Hooks: Automated Guardrails",
    [
        "Pre-commit hooks — Run before any commit:",
        "  Lint, type-check, run tests, validate formatting",
        "",
        "Post-edit hooks — Run after Claude edits files:",
        "  Run related tests, check for regressions",
        "",
        "Why enterprise customers love this:",
        "  Quality gates that enforce themselves",
        "  No human in the loop needed for standards compliance",
        "  Same rules for every developer — and for Claude",
    ],
    accent_color=GREEN,
)

# 3.3 Live Demo: Hooks in Action
make_content_slide(
    "Live Demo: Hooks in Action",
    [
        "1.  Claude makes a change that introduces a bug",
        "2.  Claude tries to commit",
        "3.  Pre-commit hook catches failing test",
        "4.  Claude reads the failure, auto-fixes the bug",
        "5.  Claude retries the commit — tests pass",
        "",
        "The agentic loop is self-correcting.",
        "Hooks make the loop reliable.",
    ],
    accent_color=GREEN,
    subtitle="Facilitator demo script",
)

# 3.4 MCP: Connecting Claude to Everything
make_content_slide(
    "MCP: Connecting Claude to Everything",
    [
        "Model Context Protocol — a standard for tool connections.",
        "",
        "MCP servers connect Claude to:",
        "  Jira / Linear — pull ticket context, update status",
        "  Slack — read threads, post updates",
        "  Datadog — query metrics, check alerts",
        "  Confluence — read docs, architecture diagrams",
        "  Internal APIs — any REST/GraphQL endpoint",
        "",
        "Claude discovers available tools dynamically.",
        "Setup pattern is the same everywhere.",
    ],
    accent_color=GREEN,
)

# 3.5 Live Demo: MCP + Jira
make_content_slide(
    "Live Demo: MCP + Jira",
    [
        "1.  Connect Claude to Jira via MCP",
        "2.  \"Pull the details for PROJ-1234\"",
        "3.  Claude reads ticket: requirements, acceptance criteria",
        "4.  Claude implements the feature, guided by ticket context",
        "5.  Claude runs tests, commits with ticket reference",
        "",
        "Developer never leaves the terminal.",
        "Context flows from tool to code to commit.",
    ],
    accent_color=GREEN,
    subtitle="Facilitator demo script",
)

# 3.6 Slash Commands: Team Workflows
make_content_slide(
    "Slash Commands: Team Workflows",
    [
        "Defined in .claude/commands/ — shared via git.",
        "",
        "Tech lead defines once, entire team inherits:",
        "",
        "  /deploy-check — Pre-deployment checklist",
        "  /review — Code review against team standards",
        "  /onboard — Set up a new developer's environment",
        "  /incident — Investigate a production issue",
        "",
        "Workflows become versioned, reviewable, repeatable.",
    ],
    accent_color=GREEN,
)

# 3.7 The Composed Workflow
make_content_slide(
    "The Composed Workflow",
    [
        "The full enterprise loop:",
        "",
        "1.  MCP pulls ticket context from Jira",
        "2.  Claude implements the feature",
        "3.  Hooks enforce quality gates (lint, test, type-check)",
        "4.  Slash command runs pre-deploy checklist",
        "5.  Code is committed with ticket reference",
        "",
        "This is the enterprise pitch:",
        "  Context → Implementation → Quality → Deployment",
        "  Automated, auditable, repeatable.",
    ],
    accent_color=GREEN,
)

# 3.8 Discussion: Your Customer's Integration
make_content_slide(
    "Discussion: Your Customer's Integration",
    [
        "What MCP servers would your customer need?",
        "",
        "What hooks would they want on day one?",
        "",
        "What slash commands would save their team the most time?",
        "",
        "How would you pitch the composed workflow?",
    ],
    accent_color=GREEN,
    subtitle="Group discussion — 10 minutes",
)


# ═════════════════════════════════════════════════════════════════════════
# DAY 4: CUSTOMER SCENARIOS (orange)
# ═════════════════════════════════════════════════════════════════════════
make_section_divider("DAY 4", "Customer Scenarios",
                     "Security, deployment, and competitive positioning", ORANGE)

# 4.1 Three Conversations, Three Playbooks
make_content_slide(
    "Three Conversations, Three Playbooks",
    [
        "Today's format: paired role-plays with facilitator coaching.",
        "",
        "1.  CISO Conversation — Security architecture and trust",
        "",
        "2.  VP Engineering Conversation — Deployment, cost, and ROI",
        "",
        "3.  Engineering Manager Conversation — Competitive positioning",
        "",
        "25 minutes each, then rotate. Real objections, real answers.",
    ],
    accent_color=ORANGE,
)

# 4.2 The Security Story
make_content_slide(
    "The Security Story",
    [
        "Layers of defense, bottom to top:",
        "",
        "OS-Level Sandboxing",
        "  Network disabled by default. Filesystem restricted.",
        "",
        "Tool-Level Permissions",
        "  Claude asks before running bash, editing files outside project.",
        "",
        "Team-Level Hooks",
        "  Pre-commit gates enforce standards automatically.",
        "",
        "Admin-Level Managed Settings",
        "  Organization-wide policies, centrally controlled.",
    ],
    accent_color=ORANGE,
    subtitle="Build confidence layer by layer",
)

# 4.3 Key Security Points
make_content_slide(
    "Key Security Points",
    [
        "Sandboxing — Network disabled by default, filesystem restricted",
        "",
        "Permissions — Claude asks before bash execution, file edits",
        "",
        "Data — No training on API inputs. SOC 2 Type II certified.",
        "",
        "Trust Center — trust.anthropic.com",
        "",
        "Key message for CISOs:",
        '  "We built the security model for the most cautious',
        '   organizations. Every layer is opt-in and auditable."',
    ],
    accent_color=ORANGE,
)

# 4.4 The Deployment Architecture
make_content_slide(
    "The Deployment Architecture",
    [
        "Deployment options:",
        "",
        "Direct API — Fastest to start, Anthropic-hosted",
        "",
        "AWS Bedrock — Data stays in your AWS account",
        "",
        "Google Vertex — Data stays in your GCP project",
        "",
        "Azure Foundry — Data stays in your Azure subscription",
        "",
        "Decision framework: data residency, existing cloud provider,",
        "  compliance requirements, procurement path.",
    ],
    accent_color=ORANGE,
)

# 4.5 Cost Math That Lands
make_content_slide(
    "Cost Math That Lands",
    [
        "Average cost: ~$6 per developer per day",
        "",
        "200 developers  x  $6/day  =  $1,200/day",
        "",
        "If each dev saves just 1 hour:",
        "  200  x  $150/hour  =  $30,000/day recovered",
        "",
        "That's a 25x ROI.",
        "",
        "Reframe the question:",
        '  "What does NOT having this cost you — in velocity,',
        '   in developer retention, in competitive speed?"',
    ],
    accent_color=ORANGE,
)

# 4.6 Competitive Positioning
make_content_slide(
    "Competitive Positioning",
    [
        "NOT \"Claude Code vs. Copilot.\"",
        "",
        "Copilot = line-level autocomplete",
        "  Excellent at what it does. Fast, lightweight, familiar.",
        "",
        "Claude Code = project-level agentic coding",
        "  Different paradigm. Plans, executes, verifies across files.",
        "",
        "They often coexist. The real question:",
        '  "What can your team do now that they couldn\'t before?"',
        "",
        "Shift from tools to capabilities.",
    ],
    accent_color=ORANGE,
)

# 4.7 Role-Play Briefing
make_content_slide(
    "Role-Play Briefing",
    [
        "Three scenarios — 25 minutes each, then rotate:",
        "",
        "1.  Nova Insurance — CISO concerned about data handling,",
        "    model access to source code, and compliance posture.",
        "",
        "2.  Atlas Manufacturing — VP Engineering evaluating deployment",
        "    options, cost model, and rollout strategy for 500 devs.",
        "",
        "3.  Prism Analytics — Engineering Manager comparing Claude Code",
        "    to existing Copilot deployment. Team is skeptical.",
        "",
        "Partner up. One plays customer, one plays Anthropic.",
    ],
    accent_color=ORANGE,
    subtitle="Paired role-plays with facilitator coaching",
)

# 4.8 Debrief Prompts
make_content_slide(
    "Debrief Prompts",
    [
        "Which objections were hardest to handle?",
        "",
        "Where were you most confident?",
        "",
        "What do you need to practice before doing this live?",
        "",
        "What would you change about your approach?",
        "",
        "What questions did you not have answers for?",
    ],
    accent_color=ORANGE,
    subtitle="Group debrief — 15 minutes",
)


# ═════════════════════════════════════════════════════════════════════════
# DAY 5: SHIP IT (green)
# ═════════════════════════════════════════════════════════════════════════
make_section_divider("DAY 5", "Ship It",
                     "Blind brief to live demo", GREEN)

# 5.1 The Capstone Format
make_content_slide(
    "The Capstone Format",
    [
        "120 minutes. Fully integrated. No safety net.",
        "",
        "Brief distribution                              0:00",
        "Analysis — understand the customer              0:15",
        "Build — create your demo with Claude Code       0:30",
        "Prep — rehearse your presentation               1:15",
        "Presentations — 10 min each                     1:25",
        "Feedback + scoring                              2:00",
        "",
        "Everything you've learned this week — applied.",
    ],
    accent_color=GREEN,
)

# 5.2 What Great Demos Look Like
make_content_slide(
    "What Great Demos Look Like",
    [
        "Start with THEIR problem, not your solution.",
        "",
        "Limit to 3-4 features. Depth over breadth.",
        "",
        "Build for reliability, not flash.",
        "  If it might break, cut it.",
        "",
        "Test twice before presenting.",
        "  Once for functionality, once for flow.",
        "",
        "Show the agentic loop — let them see Claude think.",
    ],
    accent_color=GREEN,
)

# 5.3 Presentation Structure
make_content_slide(
    "Presentation Structure",
    [
        "10 minutes total:",
        "",
        "Problem Framing  (2 min)",
        '  "Here\'s what we heard. Here\'s what matters to you."',
        "",
        "Live Demo  (6 min)",
        "  Show Claude solving their actual problem. Narrate the loop.",
        "",
        "Next Steps + Rollout Plan  (2 min)",
        "  Concrete actions. Pilot size, timeline, success metrics.",
        "",
        "Lead with empathy. Close with action.",
    ],
    accent_color=GREEN,
)

# 5.4 Scoring Rubric
make_content_slide(
    "Scoring Rubric",
    [
        "Cohort scores each presentation on four dimensions:",
        "",
        "Problem Framing",
        "  Did you understand their world?",
        "",
        "Technical Depth",
        "  Did your demo actually work?",
        "",
        "Relevance",
        "  Did it map to their specific pain points?",
        "",
        "Confidence",
        "  Did you handle questions smoothly?",
    ],
    accent_color=GREEN,
)

# 5.5 Go Time
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0), Inches(0), SLIDE_W, Inches(0.18))
shape.fill.solid()
shape.fill.fore_color.rgb = GREEN
shape.line.fill.background()

add_textbox(slide, Inches(0.75), Inches(2.5), Inches(11), Inches(1.0),
            "Go Time", font_size=48, bold=True, color=TEXT)
add_textbox(slide, Inches(0.75), Inches(3.7), Inches(11), Inches(0.6),
            "Briefs distributed. Timer starts.", font_size=24, color=MUTED)
add_textbox(slide, Inches(0.75), Inches(4.5), Inches(11), Inches(0.6),
            "Good luck.", font_size=24, color=GREEN, bold=True)

add_footer(slide)


# ═════════════════════════════════════════════════════════════════════════
# CLOSING SLIDE
# ═════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)

# Accent stripe
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0), Inches(0), SLIDE_W, Inches(0.12))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()

add_textbox(slide, Inches(0.75), Inches(2.2), Inches(11), Inches(1.0),
            "You're ready.", font_size=48, bold=True, color=TEXT)
add_textbox(slide, Inches(0.75), Inches(3.5), Inches(11), Inches(0.6),
            "Claude Code Basecamp — Complete", font_size=24, color=MUTED)

# Small accent bar
add_accent_bar(slide, ORANGE, top=Inches(4.4), width=Inches(2.0), height=Inches(0.05))

add_textbox(slide, Inches(0.75), Inches(5.0), Inches(11), Inches(0.5),
            "anthropic.com", font_size=16, color=FAINT)

add_footer(slide)


# ═════════════════════════════════════════════════════════════════════════
# SAVE
# ═════════════════════════════════════════════════════════════════════════
output_path = "/Users/alexamurray/claude-code-basecamp/public/claude-code-basecamp-slides.pptx"
prs.save(output_path)
print(f"Saved {slide_number[0]} slides to {output_path}")
